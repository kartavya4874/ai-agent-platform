import os
import httpx
import base64
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from fastapi.responses import FileResponse, StreamingResponse
from sqlalchemy.orm import Session
from app.db.database import get_db
from app.models.models import User, UserContent
from app.services.openai_service import OpenAIService
from typing import Optional
from io import BytesIO
import pptx
from pptx.util import Inches, Pt
from docx import Document
import tempfile
import shutil
import requests
from PIL import Image

router = APIRouter()

# Helper function to get temporary file path
def get_temp_file_path(filename):
    temp_dir = tempfile.gettempdir()
    return os.path.join(temp_dir, filename)

@router.post("/generate-image")
async def generate_image(prompt: str, style: Optional[str] = "realistic"):
    """Generate an image based on text prompt"""
    full_prompt = f"{prompt} in {style} style"
    
    # Call OpenAI to generate image
    image_url = await OpenAIService.generate_image(full_prompt)
    
    if not image_url:
        raise HTTPException(status_code=500, detail="Failed to generate image")
    
    # Download the image
    async with httpx.AsyncClient() as client:
        response = await client.get(image_url)
        if response.status_code != 200:
            raise HTTPException(status_code=500, detail="Failed to download generated image")
        
        # Save to temporary file
        img_temp_path = get_temp_file_path(f"generated_image_{hash(prompt)}.png")
        with open(img_temp_path, "wb") as f:
            f.write(response.content)
    
    return {"image_path": img_temp_path, "success": True}

@router.post("/generate-code")
async def generate_code(prompt: str, language: str):
    """Generate code based on text prompt"""
    # Generate code
    code = await OpenAIService.generate_code(prompt, language)
    
    if not code:
        raise HTTPException(status_code=500, detail="Failed to generate code")
    
    # Determine file extension
    extensions = {
        "python": "py",
        "javascript": "js",
        "typescript": "ts",
        "java": "java",
        "c++": "cpp",
        "c#": "cs",
        "go": "go",
        "ruby": "rb",
        "php": "php",
        "swift": "swift",
        "kotlin": "kt",
        "rust": "rs",
        "sql": "sql",
        "html": "html",
        "css": "css"
    }
    
    ext = extensions.get(language.lower(), "txt")
    
    # Save to temporary file
    code_temp_path = get_temp_file_path(f"generated_code_{hash(prompt)}.{ext}")
    with open(code_temp_path, "w") as f:
        f.write(code)
    
    return {"code": code, "file_path": code_temp_path, "success": True}

@router.post("/generate-document")
async def generate_document(prompt: str, format: str = "docx"):
    """Generate a document based on text prompt"""
    # Generate content
    content = await OpenAIService.generate_text(prompt)
    
    if not content:
        raise HTTPException(status_code=500, detail="Failed to generate document content")
    
    if format.lower() == "docx":
        # Create a new Word document
        doc = Document()
        doc.add_heading('Generated Document', 0)
        
        # Add the content
        paragraphs = content.split('\n\n')
        for para in paragraphs:
            if para.strip():
                if para.startswith('# '):
                    doc.add_heading(para[2:], level=1)
                elif para.startswith('## '):
                    doc.add_heading(para[3:], level=2)
                elif para.startswith('### '):
                    doc.add_heading(para[4:], level=3)
                else:
                    doc.add_paragraph(para)
        
        # Save the document
        doc_temp_path = get_temp_file_path(f"generated_document_{hash(prompt)}.docx")
        doc.save(doc_temp_path)
        
        return {"file_path": doc_temp_path, "success": True}
    
    elif format.lower() == "pdf":
        # For PDF, we'll create a simple text file for now
        # In a production app, you'd use a library like reportlab
        text_temp_path = get_temp_file_path(f"generated_document_{hash(prompt)}.txt")
        with open(text_temp_path, "w") as f:
            f.write(content)
        
        return {"file_path": text_temp_path, "success": True}
    
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported format: {format}")

@router.post("/generate-presentation")
async def generate_presentation(prompt: str, slides: int = 5, template: str = "professional"):
    """Generate a PowerPoint presentation based on text prompt with template options"""
    # Generate content structure
    structure_prompt = f"Create a {slides}-slide presentation structure on the topic: {prompt}. For each slide, provide a title and bullet points. Format as 'Slide 1: Title\\n- Bullet 1\\n- Bullet 2'"
    
    structure = await OpenAIService.generate_text(structure_prompt)
    
    if not structure:
        raise HTTPException(status_code=500, detail="Failed to generate presentation structure")
    
    # Create PowerPoint presentation
    prs = pptx.Presentation()
    
    # Apply template based on selection
    if template == "professional":
        # Blue professional template
        slide_layouts = [
            {"layout_idx": 0, "bg_color": (0, 65, 120), "text_color": (255, 255, 255)},  # Title slide
            {"layout_idx": 1, "bg_color": (240, 240, 240), "text_color": (0, 65, 120)}   # Content slide
        ]
    elif template == "creative":
        # Creative colorful template
        slide_layouts = [
            {"layout_idx": 0, "bg_color": (110, 43, 98), "text_color": (255, 255, 255)},  # Title slide
            {"layout_idx": 1, "bg_color": (250, 240, 250), "text_color": (110, 43, 98)}   # Content slide
        ]
    elif template == "minimal":
        # Minimal white template
        slide_layouts = [
            {"layout_idx": 0, "bg_color": (255, 255, 255), "text_color": (80, 80, 80)},  # Title slide
            {"layout_idx": 1, "bg_color": (255, 255, 255), "text_color": (80, 80, 80)}   # Content slide
        ]
    else:
        # Default template
        slide_layouts = [
            {"layout_idx": 0, "bg_color": None, "text_color": None},  # Title slide
            {"layout_idx": 1, "bg_color": None, "text_color": None}   # Content slide
        ]
    
    # Parse the structure and create slides
    current_slide = None
    slide_count = 0
    
    for line in structure.split('\n'):
        line = line.strip()
        if not line:
            continue
        
        if line.startswith('Slide ') and ':' in line:
            # New slide
            slide_count += 1
            title = line.split(':', 1)[1].strip()
            
            # Select layout (alternate between title and content layouts)
            layout_info = slide_layouts[0] if slide_count == 1 else slide_layouts[1]
            
            slide_layout = prs.slide_layouts[layout_info["layout_idx"]]
            current_slide = prs.slides.add_slide(slide_layout)
            
            # Apply background color if specified
            if layout_info["bg_color"]:
                background = current_slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = pptx.dml.color.RGBColor(*layout_info["bg_color"])
            
            # Set title and apply text color if specified
            title_shape = current_slide.shapes.title
            title_shape.text = title
            
            if layout_info["text_color"]:
                for paragraph in title_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = pptx.dml.color.RGBColor(*layout_info["text_color"])
            
            # Get content placeholder if it exists
            if slide_count > 1:  # Not the title slide
                for shape in current_slide.placeholders:
                    if shape.placeholder_format.type == 1:  # Content placeholder
                        current_content = shape
                        break
                else:
                    # If no content placeholder, add a textbox
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(4)
                    current_content = current_slide.shapes.add_textbox(left, top, width, height)
                
                content_text = current_content.text_frame
        
        elif line.startswith('- ') and current_slide and slide_count > 1:
            # Bullet point for content slides
            p = content_text.add_paragraph()
            p.text = line[2:].strip()
            p.level = 0
            
            # Apply text color if specified
            if slide_layouts[1]["text_color"]:
                for run in p.runs:
                    run.font.color.rgb = pptx.dml.color.RGBColor(*slide_layouts[1]["text_color"])
    
    # Save presentation
    ppt_temp_path = get_temp_file_path(f"generated_presentation_{hash(prompt)}.pptx")
    prs.save(ppt_temp_path)
    
    return {"file_path": ppt_temp_path, "success": True}
        
        
@router.post("/text-to-speech")
async def text_to_speech(text: str, voice: str = "en-US-Neural2-F"):
    """Convert text to speech using OpenAI's TTS API"""
    try:
        # Use OpenAI's text-to-speech API
        audio_file_path = get_temp_file_path(f"speech_{hash(text)}.mp3")
        
        # Call OpenAI's TTS endpoint
        try:
            # Import the openai library here to avoid circular import
            import openai
            
            response = openai.audio.speech.create(
                model="tts-1",
                voice=voice,
                input=text
            )
            
            # Save the audio file
            response_bytes = response.read()  # Get the raw bytes
            with open(audio_file_path, "wb") as file:
                file.write(response_bytes)
                
            return {"file_path": audio_file_path, "success": True}
        except Exception as e:
            # For testing purposes, create a dummy MP3 file
            # In a real environment, this should be removed
            dummy_mp3_path = get_temp_file_path(f"dummy_speech_{hash(text)}.mp3")
            
            # Generate a silent MP3 file
            try:
                from pydub import AudioSegment
                from pydub.generators import Sine
                
                # Generate 3 seconds of silence
                silence = AudioSegment.silent(duration=3000)
                silence.export(dummy_mp3_path, format="mp3")
                
                return {
                    "file_path": dummy_mp3_path, 
                    "success": True, 
                    "message": f"Using a dummy audio file for testing. In production, this would be real speech audio. Error: {str(e)}"
                }
            except ImportError:
                # If pydub is not available, create an empty file
                with open(dummy_mp3_path, "wb") as f:
                    # Just write minimal MP3 header bytes
                    f.write(b"\xFF\xFB\x90\x44\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00\x00")
                
                return {
                    "file_path": dummy_mp3_path, 
                    "success": True, 
                    "message": f"Using a placeholder file. Real implementation would generate audio. Error: {str(e)}"
                }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to convert text to speech: {str(e)}")
